from __future__ import annotations

from pathlib import Path

from src.pptx_vendor import ensure_pptx_vendor
from src.presentation_deck_builder import load_presentation_data
from src.presentation_editable_pptx_builder import build_editable_presentation
from src.presentation_visual_assets import build_presentation_visual_assets

ensure_pptx_vendor()

from pptx import Presentation


def test_presentation_data_matches_city_held_out_contract() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    data = load_presentation_data(repo_root)

    assert data.city_count == 30
    assert data.outer_fold_count == 5
    assert data.held_out_cities_per_fold == 6
    assert data.row_count == 71_394_894
    assert data.target_column == "hotspot_10pct"
    assert data.rf_frontier.pooled_pr_auc > data.logistic_5k.pooled_pr_auc
    assert data.partner_rf.class_1_f1_mean > data.partner_logistic.class_1_f1_mean
    assert round(data.partner_support_fraction_mean, 1) == 0.3


def test_reusable_presentation_figures_are_written() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    data = load_presentation_data(repo_root)
    assets = build_presentation_visual_assets(repo_root, data)

    assert assets.predictors_schematic_png.exists()
    assert assets.predictors_schematic_svg.exists()
    assert assets.evaluation_questions_png.exists()
    assert assets.evaluation_questions_svg.exists()
    assert assets.within_city_results_png.exists()
    assert assets.within_city_results_svg.exists()
    assert assets.transfer_results_png.exists()
    assert assets.transfer_results_svg.exists()
    assert assets.contrast_takeaway_png.exists()
    assert assets.contrast_takeaway_svg.exists()


def test_editable_pptx_has_multiple_objects_and_editable_text(tmp_path: Path) -> None:
    repo_root = Path(__file__).resolve().parents[1]
    output = tmp_path / "editable_deck.pptx"
    build_editable_presentation(repo_root=repo_root, output_path=output)

    prs = Presentation(output)
    assert len(prs.slides) == 7

    for slide in prs.slides:
        assert all(not shape.is_placeholder for shape in slide.shapes)
        assert any(getattr(shape, "has_text_frame", False) for shape in slide.shapes)

    title_text = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
    assert any("Nicholas Machado" in text for text in title_text)

    problem_slide = prs.slides[1]
    assert sum(1 for shape in problem_slide.shapes if shape.shape_type == 13) == 1
    assert any(
        "Research question + predictors" in shape.text
        for shape in problem_slide.shapes
        if hasattr(shape, "text")
    )

    design_slide = prs.slides[2]
    assert sum(1 for shape in design_slide.shapes if shape.shape_type == 13) == 1
    assert any(
        "Two evaluation questions" in shape.text
        for shape in design_slide.shapes
        if hasattr(shape, "text")
    )

    results_slide = prs.slides[3]
    assert sum(1 for shape in results_slide.shapes if shape.shape_type == 13) == 1
    assert any(
        "Within-city held-out evaluation" in shape.text
        for shape in results_slide.shapes
        if hasattr(shape, "text")
    )
    assert not any(
        "weaker" in shape.text.lower()
        for shape in results_slide.shapes
        if hasattr(shape, "text")
    )

    transfer_slide = prs.slides[4]
    assert sum(1 for shape in transfer_slide.shapes if shape.shape_type == 13) == 1
    assert any(
        "City-held-out transfer evaluation" in shape.text
        for shape in transfer_slide.shapes
        if hasattr(shape, "text")
    )

    takeaway_slide = prs.slides[5]
    assert sum(1 for shape in takeaway_slide.shapes if shape.shape_type == 13) == 1
    assert any(
        "What the contrast shows" in shape.text
        for shape in takeaway_slide.shapes
        if hasattr(shape, "text")
    )
    assert any(
        "Basic factors contain real hotspot signal" in shape.text
        for shape in takeaway_slide.shapes
        if hasattr(shape, "text")
    )

    qa_slide = prs.slides[6]
    assert sum(1 for shape in qa_slide.shapes if shape.shape_type == 13) == 0
    assert any("Questions?" in shape.text for shape in qa_slide.shapes if hasattr(shape, "text"))
    assert not any("STAT 5630" in shape.text for shape in qa_slide.shapes if hasattr(shape, "text"))

    undersized_runs: list[tuple[int, str, float | None]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    size = run.font.size.pt if run.font.size else None
                    if size is None or size < 18:
                        undersized_runs.append((slide_idx, run.text, size))
    assert undersized_runs == []
