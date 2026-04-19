from __future__ import annotations

from pathlib import Path

from src.pptx_vendor import ensure_pptx_vendor
from src.presentation_deck_builder import PresentationData, load_presentation_data
from src.presentation_visual_assets import PresentationVisualAssets, build_presentation_visual_assets

ensure_pptx_vendor()

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
SLIDE_ASSET_WIDTH = 1600
SLIDE_ASSET_HEIGHT = 900

BACKGROUND = "F6F1E7"
INK = "1E2A2F"
MUTED = "5C696D"
ACCENT = "D26B34"
ACCENT_DARK = "9A402D"
TEAL = "3A7182"
TEAL_DARK = "254D5A"
WHITE = "FFFDFA"
LIGHT = "EBE0CF"
OUTLINE = "DBC7AA"
BODY_FONT = "Arial"


def build_editable_presentation(repo_root: Path, output_path: Path) -> Path:
    data = load_presentation_data(repo_root)
    visual_assets = build_presentation_visual_assets(repo_root, data)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    blank = prs.slide_layouts[6]

    _slide_title(prs.slides.add_slide(blank))
    _slide_research_question(prs.slides.add_slide(blank), visual_assets)
    _slide_evaluation_questions(prs.slides.add_slide(blank), visual_assets)
    _slide_within_city_results(prs.slides.add_slide(blank), data, visual_assets)
    _slide_transfer_results(prs.slides.add_slide(blank), visual_assets)
    _slide_contrast_takeaway(prs.slides.add_slide(blank), visual_assets)
    _slide_qa(prs.slides.add_slide(blank))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def _x(px: float) -> Inches:
    return Inches(px / SLIDE_ASSET_WIDTH * SLIDE_WIDTH_IN)


def _y(py: float) -> Inches:
    return Inches(py / SLIDE_ASSET_HEIGHT * SLIDE_HEIGHT_IN)


def _w(px: float) -> Inches:
    return Inches(px / SLIDE_ASSET_WIDTH * SLIDE_WIDTH_IN)


def _h(py: float) -> Inches:
    return Inches(py / SLIDE_ASSET_HEIGHT * SLIDE_HEIGHT_IN)


def _set_background(slide, color: str = BACKGROUND) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_oval(slide, x, y, w, h, color: str, transparency: float) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, _x(x), _y(y), _w(w), _h(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.fill.transparency = transparency
    shape.line.fill.background()


def _add_round_rect(
    slide,
    x,
    y,
    w,
    h,
    fill_color: str,
    line_color: str | None = None,
    line_width_pt: float = 1.0,
) -> object:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, _x(x), _y(y), _w(w), _h(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line_color)
        shape.line.width = Pt(line_width_pt)
    return shape


def _add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    font_size: float,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = BODY_FONT,
    margin_pt: float = 4.0,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    auto_fit: bool = False,
) -> object:
    box = slide.shapes.add_textbox(_x(x), _y(y), _w(w), _h(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(margin_pt)
    tf.margin_right = Pt(margin_pt)
    tf.margin_top = Pt(margin_pt)
    tf.margin_bottom = Pt(margin_pt)
    tf.vertical_anchor = valign
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
    return box


def _add_picture(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), _x(x), _y(y), _w(w), _h(h))


def _slide_title(slide) -> None:
    _set_background(slide)
    _add_oval(slide, 1180, -60, 360, 360, ACCENT, 0.90)
    _add_oval(slide, 1280, 560, 220, 220, TEAL, 0.93)

    _add_textbox(
        slide,
        84,
        180,
        980,
        250,
        "Cross-City Urban Heat\nHotspot Prediction",
        39,
        INK,
        True,
        auto_fit=True,
    )
    _add_textbox(
        slide,
        88,
        446,
        920,
        48,
        "Two ways to evaluate hotspot prediction",
        23,
        ACCENT_DARK,
        True,
        auto_fit=True,
    )
    _add_textbox(
        slide,
        88,
        506,
        700,
        44,
        "Max Clements | Nicholas Machado",
        23,
        TEAL_DARK,
        True,
        auto_fit=True,
    )

    _add_round_rect(slide, 72, 790, 1360, 60, LIGHT)
    _add_textbox(
        slide,
        96,
        804,
        1310,
        28,
        "STAT 5630 Final Project | April 2026",
        18,
        TEAL_DARK,
        False,
        auto_fit=True,
    )


def _slide_research_question(slide, visual_assets: PresentationVisualAssets) -> None:
    _set_background(slide)
    _add_textbox(
        slide,
        58,
        42,
        1080,
        44,
        "Research question + predictors",
        25,
        INK,
        True,
        auto_fit=True,
    )
    _add_round_rect(slide, 42, 112, 1516, 548, WHITE, OUTLINE, 1.2)
    _add_picture(slide, visual_assets.predictors_schematic_png, 64, 132, 1472, 504)
    _add_round_rect(slide, 118, 692, 1360, 64, LIGHT, OUTLINE, 0.8)
    _add_textbox(
        slide,
        150,
        706,
        1290,
        34,
        "The target is city-relative: the hottest 10% of 30 m cells within each city.",
        18,
        TEAL_DARK,
        True,
        PP_ALIGN.CENTER,
        auto_fit=True,
    )


def _slide_evaluation_questions(slide, visual_assets: PresentationVisualAssets) -> None:
    _set_background(slide)
    _add_textbox(slide, 58, 42, 980, 44, "Two evaluation questions", 25, INK, True, auto_fit=True)
    _add_round_rect(slide, 42, 112, 1516, 578, WHITE, OUTLINE, 1.2)
    _add_picture(slide, visual_assets.evaluation_questions_png, 64, 130, 1472, 534)
    _add_textbox(
        slide,
        130,
        724,
        1340,
        36,
        "The split design changes the use case being measured.",
        18,
        TEAL_DARK,
        True,
        PP_ALIGN.CENTER,
        auto_fit=True,
    )


def _slide_within_city_results(
    slide, data: PresentationData, visual_assets: PresentationVisualAssets
) -> None:
    _set_background(slide)
    _add_textbox(
        slide,
        52,
        36,
        1100,
        60,
        "Within-city held-out evaluation",
        25,
        INK,
        True,
        auto_fit=True,
    )

    _add_round_rect(slide, 44, 112, 1512, 540, WHITE, OUTLINE, 1.2)
    _add_picture(slide, visual_assets.within_city_results_png, 84, 132, 1432, 496)

    _add_round_rect(slide, 86, 684, 1428, 74, LIGHT, OUTLINE, 0.8)
    _add_textbox(
        slide,
        126,
        698,
        1348,
        38,
        f"We treat these results as within-city / row-level evidence because support counts are about {data.partner_support_fraction_mean:.0%} of each city's rows.",
        18,
        MUTED,
        False,
        PP_ALIGN.CENTER,
        auto_fit=True,
    )


def _slide_transfer_results(slide, visual_assets: PresentationVisualAssets) -> None:
    _set_background(slide)
    _add_textbox(
        slide,
        52,
        36,
        1320,
        46,
        "City-held-out transfer evaluation",
        25,
        INK,
        True,
        auto_fit=True,
    )
    _add_round_rect(slide, 44, 112, 1512, 540, WHITE, OUTLINE, 1.2)
    _add_picture(slide, visual_assets.transfer_results_png, 84, 132, 1432, 496)
    _add_round_rect(slide, 86, 684, 1428, 74, LIGHT, OUTLINE, 0.8)
    _add_textbox(
        slide,
        126,
        698,
        1348,
        38,
        "To measure transfer to new cities, whole cities are withheld during model fitting and tuning.",
        18,
        MUTED,
        False,
        PP_ALIGN.CENTER,
        auto_fit=True,
    )


def _slide_contrast_takeaway(slide, visual_assets: PresentationVisualAssets) -> None:
    _set_background(slide)
    _add_textbox(
        slide,
        52,
        42,
        1460,
        72,
        "What the contrast shows",
        25,
        INK,
        True,
        auto_fit=True,
    )
    _add_round_rect(slide, 44, 124, 1512, 564, WHITE, OUTLINE, 1.2)
    _add_picture(slide, visual_assets.contrast_takeaway_png, 68, 144, 1464, 518)
    _add_textbox(
        slide,
        150,
        724,
        1300,
        36,
        "Basic factors contain real hotspot signal, but transfer to unseen cities is the hard part.",
        19,
        TEAL_DARK,
        True,
        PP_ALIGN.CENTER,
        auto_fit=True,
    )


def _slide_qa(slide) -> None:
    _set_background(slide)
    _add_oval(slide, 1220, 40, 320, 320, ACCENT, 0.9)
    _add_oval(slide, 1280, 560, 220, 220, TEAL, 0.93)

    _add_textbox(slide, 78, 246, 560, 88, "Questions?", 42, INK, True, auto_fit=True)
